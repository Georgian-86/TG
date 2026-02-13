"""
Presentation Agent
Generates PowerPoint and PDF files from lesson content
Async class-based implementation with BaseAgent inheritance
"""
from typing import Dict, Any, List, Optional
import asyncio
import os
import re
import textwrap
from pathlib import Path
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF

from app.agents.base import BaseAgent
from app.config import settings

logger = logging.getLogger(__name__)

# =========================================================
# CONFIGURATION
# =========================================================

OUTPUT_DIR = Path("outputs")
MASCOT_PATH = Path("assets/TechGenieMascot.png")
FONT_REGULAR = Path("assets/fonts/DejaVuSans.ttf")
FONT_BOLD = Path("assets/fonts/DejaVuSans-Bold.ttf")


class PresentationAgent(BaseAgent):
    """Generates PowerPoint and PDF presentation files from lesson content"""
    
    def __init__(self):
        super().__init__()
        self.output_dir = OUTPUT_DIR
        self.mascot_path = MASCOT_PATH
        self.font_regular = FONT_REGULAR
        self.font_bold = FONT_BOLD
        
    def _ensure_dirs(self):
        """Create output directory if it doesn't exist"""
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
    def _clean_text(self, text: str) -> str:
        """Normalize text safely for PPT/PDF (handles smart quotes)"""
        if not text:
            return ""
        
        # Handle dict content (new structure)
        if isinstance(text, dict):
            # Try to extract text from dict
            if "text" in text:
                text = text["text"]
            elif "content" in text:
                text = text["content"]
            elif "description" in text:
                text = text["description"]
            else:
                # Convert dict to string representation
                text = str(text)
        
        # Convert to string if not already
        text = str(text)
            
        # Replace smart quotes and other common unicode characters
        replacements = {
            '\u2018': "'", '\u2019': "'",  # Smart single quotes
            '\u201c': '"', '\u201d': '"',  # Smart double quotes
            '\u2013': '-', '\u2014': '-',  # En/Em dashes
            '\u2026': '...',               # Ellipsis
        }
        for char, replacement in replacements.items():
            text = text.replace(char, replacement)
            
        # If fonts aren't available, force ASCII to prevent crashes
        if not self.font_regular.exists():
            text = text.encode('ascii', 'ignore').decode('ascii')
            
        text = re.sub(r"\s+", " ", text)
        return text.strip()
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences for bullet points"""
        text = self._clean_text(text)
        parts = re.split(r"(?<=[.!?])\s+", text)
        return [p for p in parts if len(p) > 25]
    
    def _summarize_to_bullets(self, content, max_bullets: int = 6) -> List[str]:
        """Convert text or dict to bullet points (max 6 by default)"""
        # Handle different content structures
        if isinstance(content, dict):
            # NEW: Handle nested subsections structure from Content Agent
            # Example: {"basic_syntax": "text...", "data_types": "text...", ...}
            
            # First check if it's a standard key-value structure
            if "text" in content:
                text = content["text"]
            elif "content" in content:
                text = content["content"]
            elif "description" in content:
                text = content["description"]
            else:
                # It's a subsections dict - extract all subsection texts
                subsection_texts = []
                for key, value in content.items():
                    if isinstance(value, str) and len(value) > 20:  # Extract meaningful text content
                        subsection_texts.append(value)
                    elif isinstance(value, dict):
                        # Handle nested dicts recursively
                        subsection_texts.append(str(value.get("text", value.get("content", value.get("description", "")))))
                
                # Join all subsection texts
                text = " ".join(subsection_texts)
        elif isinstance(content, list):
            # Join list items
            text = " ".join(str(item) for item in content)
        else:
            text = str(content) if content else ""
        
        sentences = self._split_into_sentences(text)
        return sentences[:max_bullets]
    
    def _extract_all_content_bullets(self, content) -> List[str]:
        """Extract ALL content from nested subsections as bullet points (no limit)"""
        bullets = []
        
        if isinstance(content, dict):
            # Handle nested subsections structure
            if "text" in content or "content" in content or "description" in content:
                # Standard dict with known keys
                text = content.get("text", content.get("content", content.get("description", "")))
                bullets.extend(self._split_into_sentences(text))
            else:
                # It's a subsections dict - extract each subsection
                for subsection_name, subsection_content in content.items():
                    if isinstance(subsection_content, str) and len(subsection_content) > 20:
                        # Split each subsection into sentences
                        sentences = self._split_into_sentences(subsection_content)
                        bullets.extend(sentences)
                    elif isinstance(subsection_content, dict):
                        # Handle nested dicts
                        nested_text = subsection_content.get("text", subsection_content.get("content", subsection_content.get("description", "")))
                        if nested_text:
                            bullets.extend(self._split_into_sentences(nested_text))
        elif isinstance(content, list):
            # Handle list of items
            for item in content:
                if isinstance(item, str):
                    bullets.extend(self._split_into_sentences(item))
                elif isinstance(item, dict):
                    text = item.get("text", item.get("content", item.get("description", "")))
                    if text:
                        bullets.extend(self._split_into_sentences(text))
        else:
            # Plain string
            text = str(content) if content else ""
            bullets.extend(self._split_into_sentences(text))
        
        # Filter out very short bullets (dynamic pagination handles length)
        return [b for b in bullets if len(b) > 25]
    
    def _add_footer(self, slide):
        """Add branding footer to slide"""
        # Footer Bar (Background)
        left = Inches(0)
        top = Inches(7.0)
        width = Inches(10)
        height = Inches(0.5)
        
        # Add a subtle footer line/bar
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE (using int to avoid import)
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb("F3F4F6") # Light gray
        shape.line.fill.background() # No border
        
        # Footer Text
        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1), Inches(9), Inches(0.4)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "Powered by TeachGenie.ai"
        p.font.size = Pt(12)
        p.font.color.rgb = self._hex_to_rgb("4B5563") # Gray-600
        p.alignment = PP_ALIGN.RIGHT
        
        # Add small mascot logo if available
        if self.mascot_path.exists():
            slide.shapes.add_picture(
                str(self.mascot_path),
                left=Inches(0.2),
                top=Inches(7.05),
                height=Inches(0.4)
            )

    def _hex_to_rgb(self, hex_color):
        """Convert hex string to RGB object"""
        from pptx.dml.color import RGBColor
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def _add_border(self, slide):
        """Add a decorative border to the slide"""
        # A4ish aspect ratio usually, but screens are 16:9 (10x5.625) or 4:3 (10x7.5)
        # Standard PPTX is 16:9 (10 inches x 5.625 inches) usually, or 4:3.
        # Let's assume standard 4:3 (10x7.5) which is default in python-pptx unless changed.
        
        # Border dimensions (leaving a small margin)
        left = Inches(0.3)
        top = Inches(0.3)
        width = Inches(9.4)
        height = Inches(6.9)
        
        # 1 = MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(1, left, top, width, height)
        
        # No fill
        shape.fill.background()
        
        # Line style
        line = shape.line
        line.color.rgb = self._hex_to_rgb("4F46E5") # Indigo-600 (TeachGenie Brand Color?)
        line.width = Pt(3)

    def _extract_content_with_subsections(self, content: Any) -> List[Dict[str, Any]]:
        """
        Extract content preserving subsection structure.
        Returns list of {"subtitle": str, "bullets": List[str]}
        """
        results = []
        
        if isinstance(content, dict):
            # Check for standard keys first
            if "text" in content:
                # Treated as main content without subtitle
                bullets = self._split_into_sentences(content["text"])
                if bullets: results.append({"subtitle": None, "bullets": bullets})
            elif "content" in content:
                bullets = self._split_into_sentences(content["content"])
                if bullets: results.append({"subtitle": None, "bullets": bullets})
            elif "description" in content:
                bullets = self._split_into_sentences(content["description"])
                if bullets: results.append({"subtitle": None, "bullets": bullets})
            else:
                # It's likely a set of subsections
                for key, value in content.items():
                    # Format key as subtitle (e.g., "core_concepts" -> "Core Concepts")
                    subtitle = key.replace("_", " ").title()
                    bullets = []
                    
                    if isinstance(value, str):
                        bullets = self._split_into_sentences(value)
                    elif isinstance(value, dict):
                        text = value.get("text", value.get("content", value.get("description", "")))
                        bullets = self._split_into_sentences(text)
                    elif isinstance(value, list):
                        for item in value:
                            bullets.extend(self._split_into_sentences(str(item)))
                            
                    if bullets:
                        results.append({"subtitle": subtitle, "bullets": bullets})
        
        elif isinstance(content, list):
             bullets = []
             for item in content:
                 bullets.extend(self._split_into_sentences(str(item)))
             if bullets:
                 results.append({"subtitle": None, "bullets": bullets})
                 
        elif isinstance(content, str):
            bullets = self._split_into_sentences(content)
            if bullets:
                results.append({"subtitle": None, "bullets": bullets})
                
        return results

    async def _build_ppt(
        self,
        topic: str,
        level: str,
        duration: int,
        sections: List[Dict[str, Any]],
        takeaways: Optional[List[str]] = None,
        quiz: Optional[Dict[str, Any]] = None
    ) -> Path:
        """Build PowerPoint presentation"""
        logger.info(f"Building PPT for: {topic}")
        try:
            self._ensure_dirs()
            ppt_path = await asyncio.to_thread(
                self._generate_ppt_sync,
                topic, level, duration, sections, takeaways, quiz
            )
            return ppt_path
        except Exception as e:
            logger.error(f"PPT generation failed: {e}")
            raise
    
    def _generate_ppt_sync(
        self,
        topic: str,
        level: str,
        duration: int,
        sections: List[Dict[str, Any]],
        takeaways: Optional[List[str]],
        quiz: Optional[Dict[str, Any]]
    ) -> Path:
        """Synchronous PPT generation (called in thread executor)"""
        prs = Presentation()
        
        # ---------- TITLE SLIDE ----------
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_border(slide) # Add border
        
        # Add mascot at center if available
        if self.mascot_path.exists():
            slide.shapes.add_picture(
                str(self.mascot_path),
                left=Inches(4.0),
                top=Inches(1.5),
                height=Inches(2.0)
            )
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(2))
        tf = title_box.text_frame
        tf.clear()
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        clean_level = str(level).replace("LessonLevel.", "").replace("_", " ").title()
        p2 = tf.add_paragraph()
        p2.text = f"{clean_level} Level  |  {duration} Minutes"
        p2.font.size = Pt(20)
        p2.font.name = "Arial"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = "Generated by TeachGenie.ai"
        p3.font.size = Pt(14)
        p3.font.color.rgb = self._hex_to_rgb("6B7280")
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(30)
        
        self._add_footer(slide)
        
        # ---------- CONTENT SLIDES ----------
        for section in sections:
            section_title = section.get("title", "Untitled Section")
            content = section.get("content", "")
            
            # Extract structured content (subsections)
            structured_content = self._extract_content_with_subsections(content)
            
            if not structured_content:
                continue

            # Iterate through subsections
            for subsection in structured_content:
                subtitle = subsection["subtitle"]
                bullets = subsection["bullets"]
                
                if not bullets: continue

                # Pagination Logic
                SLIDE_HEIGHT_INCHES = 7.5
                USABLE_HEIGHT_INCHES = 5.0
                LINE_HEIGHT_INCHES = 0.35
                CHARS_PER_LINE = 80
                
                def estimate_bullet_height(bullet_text):
                    num_lines = max(1, (len(bullet_text) + CHARS_PER_LINE - 1) // CHARS_PER_LINE)
                    return num_lines * LINE_HEIGHT_INCHES + 0.15
                
                slide_groups = []
                current_slide_bullets = []
                current_height = 0
                
                for bullet in bullets:
                    bullet_height = estimate_bullet_height(bullet)
                    if current_height + bullet_height > USABLE_HEIGHT_INCHES and current_slide_bullets:
                        slide_groups.append(current_slide_bullets)
                        current_slide_bullets = [bullet]
                        current_height = bullet_height
                    else:
                        current_slide_bullets.append(bullet)
                        current_height += bullet_height
                
                if current_slide_bullets:
                    slide_groups.append(current_slide_bullets)
                
                # Create slides
                for page_num, slide_bullets in enumerate(slide_groups):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    self._add_border(slide) # Add border
                    
                    # Title: "Section: Subsection"
                    full_title = section_title
                    if subtitle:
                        full_title += f": {subtitle}"
                    
                    if len(slide_groups) > 1:
                        full_title += f" ({page_num + 1}/{len(slide_groups)})"
                    
                    slide.shapes.title.text = full_title
                    
                    # Adjust title font size if too long
                    if len(full_title) > 50:
                         slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
                    
                    # Body
                    body = slide.shapes.placeholders[1].text_frame
                    body.clear()
                    
                    for i, bullet in enumerate(slide_bullets):
                        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(18)
                        p.level = 0
                        p.space_after = Pt(10)
                    
                    self._add_footer(slide)
        
        # ---------- KEY TAKEAWAYS SLIDE ----------
        if takeaways:
            if not isinstance(takeaways, list): takeaways = []
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            self._add_border(slide)
            slide.shapes.title.text = "Key Takeaways"
            
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            
            for i, takeaway in enumerate(takeaways[:6]):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                if isinstance(takeaway, dict):
                    title = takeaway.get("title", "Key Idea")
                    desc = takeaway.get("description", "")
                    p.text = f"{title}: {desc}"
                else:
                    p.text = str(takeaway)
                p.font.size = Pt(18)
                p.space_after = Pt(10)
                
            self._add_footer(slide)

        # ---------- QUIZ SLIDES ----------
        if quiz and isinstance(quiz, dict) and "questions" in quiz:
            questions = quiz.get("questions", [])
            if questions:
                # Quiz Section Title Slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_border(slide)
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
                p = title_box.text_frame.paragraphs[0]
                p.text = "Knowledge Check\nScenario-Based Assessment"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(32)
                p.font.bold = True
                self._add_footer(slide)

                # Question Slides
                for idx, q in enumerate(questions):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    self._add_border(slide)
                    
                    # Scenario as Title
                    scenario = q.get("scenario", "")
                    question_text = q.get("question", "")
                    
                    title = slide.shapes.title
                    title.text = f"Scenario {idx + 1}"
                    title.text_frame.paragraphs[0].font.size = Pt(24)
                    
                    body = slide.shapes.placeholders[1].text_frame
                    body.clear()
                    
                    # Scenario
                    p_scen = body.paragraphs[0]
                    p_scen.text = scenario
                    p_scen.font.size = Pt(16)
                    p_scen.font.italic = True
                    p_scen.space_after = Pt(12)
                    
                    # Question
                    p_q = body.add_paragraph()
                    p_q.text = question_text
                    p_q.font.size = Pt(18)
                    p_q.font.bold = True
                    p_q.space_after = Pt(12)
                    
                    # Options
                    options = q.get("options", {})
                    for opt_key in ["A", "B", "C", "D"]:
                        if opt_key in options:
                            p_opt = body.add_paragraph()
                            p_opt.text = f"{opt_key}) {options[opt_key]}"
                            p_opt.font.size = Pt(16)
                            p_opt.level = 1
                    
                    self._add_footer(slide)

                # Answer Key Slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                self._add_border(slide)
                slide.shapes.title.text = "Answer Key"
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                
                for idx, q in enumerate(questions):
                    p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                    ans = q.get("correct_option", "?")
                    exp = q.get("explanation", "")
                    p.text = f"Q{idx + 1}: {ans} - {exp}"
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)
                
                self._add_footer(slide)

        # Save
        safe_filename = topic.replace(" ", "_").replace("/", "_")
        ppt_path = self.output_dir / f"TG-{safe_filename}.pptx"
        prs.save(str(ppt_path))
        return ppt_path
    
    async def _build_pdf(
        self,
        topic: str,
        sections: List[Dict[str, Any]],
        takeaways: Optional[List[str]] = None,
        quiz: Optional[Dict[str, Any]] = None
    ) -> Path:
        """
        Build PDF document
        
        Args:
            topic: Lesson topic
            sections: List of lesson sections
            takeaways: Optional list of key takeaways
            quiz: Optional quiz data
            
        Returns:
            Path to generated PDF file
        """
        logger.info(f"Building PDF for: {topic}")
        
        try:
            self._ensure_dirs()
            from app.utils.pdf_generator import generate_pdf_logic
            
            # Using synchronous execution via thread pool
            logger.info("Generating PDF locally via thread pool...")
            pdf_path_str = await asyncio.to_thread(
                generate_pdf_logic, 
                topic, sections, takeaways, quiz
            )
            
            pdf_path = Path(pdf_path_str)
            logger.info(f"PDF generated successfully: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            logger.error(f"PDF generation failed: {e}")
            raise
    
    # _generate_pdf_sync removed (moved to Celery task)
    
    async def run(
        self,
        topic: str,
        level: str,
        duration: int,
        sections: List[Dict[str, Any]],
        takeaways: Optional[List[str]] = None,
        quiz: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Generate both PPT and PDF presentations
        
        Args:
            topic: Lesson topic
            level: Educational level
            duration: Lesson duration in minutes
            sections: Lesson sections with content
            takeaways: Optional key takeaways
            quiz: Optional quiz data
            
        Returns:
            Dictionary with ppt_path and pdf_path
        """
        logger.info(f"Generating presentations for: {topic}")
        
        try:
            # Generate PPT and PDF in parallel for speed
            ppt_path, pdf_path = await asyncio.gather(
                self._build_ppt(topic, level, duration, sections, takeaways, quiz),
                self._build_pdf(topic, sections, takeaways, quiz)
            )
            
            return {
                "ppt_path": str(ppt_path),
                "pdf_path": str(pdf_path)
            }
            
        except Exception as e:
            logger.error(f"Presentation generation failed: {e}")
            raise
